#!/usr/bin/env python3
"""Build a Chinese PPTX deck for GeoFeedback-GS formal four-group results."""

import csv
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Cm, Pt


ROOT = Path(__file__).resolve().parents[1]
RESULT_ROOT = ROOT / "server_results_review" / "paper_results_formal_4groups_v2"
EVIDENCE_ROOT = ROOT / "server_results_review" / "paper_evidence_formal_4groups_v2"
OUT_DIR = ROOT / "outputs" / "ppt_zh_geoguardgs"
OUT_PPTX = OUT_DIR / "GeoFeedback-GS_formal_results_zh.pptx"
OUT_NOTES = OUT_DIR / "GeoFeedback-GS_formal_results_zh_speaker_notes.md"
OUT_QA = OUT_DIR / "GeoFeedback-GS_formal_results_zh_qa.json"
TEMPLATE_GLOB = "*PPT.pptx"


W, H = Cm(33.867), Cm(19.05)
BG = RGBColor(248, 249, 251)
INK = RGBColor(31, 36, 45)
MUTED = RGBColor(98, 108, 122)
BLUE = RGBColor(39, 98, 180)
CYAN = RGBColor(39, 142, 164)
GREEN = RGBColor(46, 139, 87)
ORANGE = RGBColor(217, 128, 50)
RED = RGBColor(185, 72, 66)
LINE = RGBColor(218, 224, 232)
WHITE = RGBColor(255, 255, 255)


def read_csv(path):
    with path.open("r", encoding="utf-8-sig", newline="") as f:
        return list(csv.DictReader(f))


def rgb(hex_color):
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def set_text(box, text, size=18, bold=False, color=INK, align=None):
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Cm(1.25), Cm(0.65), Cm(24), Cm(1.0))
    set_text(box, title, size=25, bold=True, color=INK)
    if subtitle:
        sub = slide.shapes.add_textbox(Cm(1.28), Cm(1.68), Cm(25), Cm(0.55))
        set_text(sub, subtitle, size=9.5, color=MUTED)
    line = slide.shapes.add_shape(1, Cm(1.25), Cm(2.38), Cm(31.3), Cm(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.color.rgb = LINE


def add_footer(slide, idx):
    box = slide.shapes.add_textbox(Cm(28.7), Cm(18.35), Cm(3.7), Cm(0.35))
    set_text(box, f"GeoFeedback-GS | {idx:02d}", size=7.5, color=MUTED, align=PP_ALIGN.RIGHT)


def add_bullets(slide, x, y, w, h, items, size=14, color=INK, gap=0.08):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(gap)
    return box


def add_card(slide, x, y, w, h, title, body, accent=BLUE):
    shp = slide.shapes.add_shape(5, Cm(x), Cm(y), Cm(w), Cm(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = LINE
    bar = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(0.08), Cm(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.color.rgb = accent
    t = slide.shapes.add_textbox(Cm(x + 0.35), Cm(y + 0.22), Cm(w - 0.55), Cm(0.55))
    set_text(t, title, size=13, bold=True, color=accent)
    b = slide.shapes.add_textbox(Cm(x + 0.35), Cm(y + 0.92), Cm(w - 0.55), Cm(h - 1.05))
    set_text(b, body, size=10.5, color=INK)


def add_image(slide, path, x, y, w=None, h=None):
    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(path)
    if w and h:
        return slide.shapes.add_picture(str(path), Cm(x), Cm(y), width=Cm(w), height=Cm(h))
    if w:
        return slide.shapes.add_picture(str(path), Cm(x), Cm(y), width=Cm(w))
    if h:
        return slide.shapes.add_picture(str(path), Cm(x), Cm(y), height=Cm(h))
    return slide.shapes.add_picture(str(path), Cm(x), Cm(y))


def add_table(slide, x, y, w, h, rows, columns, font_size=8.5):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(columns), Cm(x), Cm(y), Cm(w), Cm(h))
    table = table_shape.table
    for j, col in enumerate(columns):
        cell = table.cell(0, j)
        cell.text = col[1]
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb("E8EEF8")
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = "Microsoft YaHei"
                r.font.bold = True
                r.font.size = Pt(font_size)
                r.font.color.rgb = INK
    for i, row in enumerate(rows, start=1):
        for j, (key, _) in enumerate(columns):
            cell = table.cell(i, j)
            cell.text = str(row.get(key, ""))
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 else rgb("F5F7FA")
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = "Microsoft YaHei"
                    r.font.size = Pt(font_size)
                    r.font.color.rgb = INK
    return table_shape


def base_slide(prs, idx, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_title(slide, title, subtitle)
    add_footer(slide, idx)
    return slide


def fmt(v, n=2):
    if v in ("", None):
        return ""
    try:
        return f"{float(v):.{n}f}"
    except Exception:
        return str(v)


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    main_rows = read_csv(RESULT_ROOT / "tables" / "table_main_results.csv")
    audit_rows = read_csv(RESULT_ROOT / "tables" / "table_no_lidar_audit.csv")
    feedback_rows = read_csv(RESULT_ROOT / "tables" / "table_feedback_summary.csv")

    # Compact display rows.
    metric_rows = []
    for r in main_rows:
        metric_rows.append({
            "label": r["label"],
            "psnr": fmt(r["psnr_mean"]),
            "ssim": fmt(r["ssim_mean"], 3),
            "lpips": fmt(r["lpips_mean"], 3),
            "l1": fmt(r["l1_mean"], 3),
            "lidar": "是" if r["uses_lidar_init"] == "True" else "否",
            "fb": f"{r['feedback_valid']}/{r['feedback_total']}" if r.get("feedback_total") else "-",
        })

    template_candidates = [
        p for p in ROOT.glob(TEMPLATE_GLOB)
        if p.name != OUT_PPTX.name and p.is_file()
    ]
    template_path = template_candidates[0] if template_candidates else None
    if template_path:
        prs = Presentation(str(template_path))
        # Keep the template theme, masters, and layouts, but remove existing
        # example slides so the generated deck starts clean.
        slide_id_list = prs.slides._sldIdLst
        for slide_id in list(slide_id_list):
            r_id = slide_id.rId
            prs.part.drop_rel(r_id)
            slide_id_list.remove(slide_id)
    else:
        prs = Presentation()
        prs.slide_width = W
        prs.slide_height = H
    notes = []

    # 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("F7F9FC")
    title = slide.shapes.add_textbox(Cm(1.5), Cm(3.1), Cm(25), Cm(1.8))
    set_text(title, "GeoFeedback-GS：基于责任高斯反馈的低 LiDAR 依赖动态街景高斯重建", size=28, bold=True)
    sub = slide.shapes.add_textbox(Cm(1.55), Cm(5.25), Cm(23), Cm(1.0))
    set_text(sub, "正式四组实验结果汇报 | StreetGS / DA3-only / DA3+Feedback / Pure-Vision DA3+Feedback", size=13, color=MUTED)
    add_card(slide, 1.6, 7.2, 9.0, 3.8, "核心结论", "PV-C 在无 LiDAR 初始化和无 LiDAR 监督条件下，保持了与 LiDAR 初始化 baseline 接近的 held-out RGB 指标。", GREEN)
    add_card(slide, 11.4, 7.2, 9.0, 3.8, "证据边界", "当前证据支持单场景、四组正式实验；不能直接声称跨场景泛化或几何精度显著提升。", ORANGE)
    add_card(slide, 21.2, 7.2, 9.0, 3.8, "输出材料", "已生成主表、审计表、PSNR/L1 曲线、90 张同视角 RGB/depth 对比图和论文初稿。", BLUE)
    add_footer(slide, 1)
    notes.append("1. 开场说明：本 PPT 汇报的是当前正式四组实验能支撑的论文论点和证据边界。")

    # 2
    slide = base_slide(prs, 2, "研究问题：LiDAR 依赖到底发生在哪个阶段？")
    add_bullets(slide, 1.6, 3.05, 13.5, 5.0, [
        "Street Gaussian 类方法常依赖 LiDAR 作为几何先验或训练监督。",
        "“无 LiDAR 监督”不等于“完全无 LiDAR”：初始化仍可能使用 LiDAR pointcloud。",
        "车辆 object branch 必须保留，否则自动驾驶街景主实验不完整。",
    ], size=15)
    add_card(slide, 17.0, 3.0, 13.7, 3.0, "本文关注", "把 LiDAR 初始化、LiDAR 训练监督、DA3 结构信号、周期反馈拆开评估。", BLUE)
    add_card(slide, 17.0, 6.55, 13.7, 3.0, "可写论文问题", "DA3 与 feedback 能否在减少 LiDAR 依赖时维持 object-aware 重建质量？", CYAN)
    add_card(slide, 17.0, 10.1, 13.7, 3.0, "不能过度声明", "当前结果不是“全面超过 StreetGS”，也不是“几何显著改善”的最终证明。", ORANGE)
    notes.append("2. 重点解释 LiDAR 初始化和 LiDAR 监督的区别。")

    # 3
    slide = base_slide(prs, 3, "正式四组实验设计", "held-out test split: 245 views；object branch 全部保留")
    group_rows = [
        {"g": "A", "name": "StreetGS baseline", "init": "LiDAR", "sup": "LiDAR", "da3": "否", "fb": "否", "role": "上界/原版基线"},
        {"g": "B", "name": "DA3-only", "init": "LiDAR", "sup": "无 LiDAR", "da3": "是", "fb": "否", "role": "替代训练监督"},
        {"g": "C", "name": "DA3+Feedback", "init": "LiDAR", "sup": "无 LiDAR", "da3": "是", "fb": "是", "role": "主方法：反馈链路"},
        {"g": "PV-C", "name": "Pure-Vision DA3+Feedback", "init": "COLMAP + random-box", "sup": "无 LiDAR", "da3": "是", "fb": "是", "role": "纯视觉加分项"},
    ]
    add_table(slide, 1.4, 3.0, 30.8, 7.1, group_rows, [
        ("g", "组别"), ("name", "配置"), ("init", "初始化"), ("sup", "训练监督"), ("da3", "DA3"), ("fb", "反馈"), ("role", "论文角色")
    ], font_size=9.5)
    add_card(slide, 1.6, 11.2, 14.8, 3.2, "评估协议", "final evaluation 只采用 held-out test，避免 sampled periodic eval 被误作主结果。", GREEN)
    add_card(slide, 17.1, 11.2, 14.8, 3.2, "可审计边界", "初始化 manifest、feedback manifest、安全审计表用于证明 LiDAR 使用边界。", BLUE)
    notes.append("3. 这一页是论文实验设计的核心。")

    # 4
    slide = base_slide(prs, 4, "GeoFeedback-GS 方法概览：DA3 信号 + 周期反馈")
    x0, y0 = 1.6, 4.0
    steps = [
        ("输入图像/相机", BLUE),
        ("COLMAP 或 LiDAR 初始化", CYAN),
        ("Street Gaussian 训练", BLUE),
        ("DA3 depth/edge/risk", GREEN),
        ("Feedback controller", ORANGE),
        ("softpatch mask", RED),
    ]
    for i, (txt, col) in enumerate(steps):
        x = x0 + i * 5.05
        shp = slide.shapes.add_shape(5, Cm(x), Cm(y0), Cm(4.25), Cm(1.55))
        shp.fill.solid()
        shp.fill.fore_color.rgb = WHITE
        shp.line.color.rgb = col
        set_text(shp.text_frame.paragraphs[0]._parent, txt, size=10.5, bold=True, color=col, align=PP_ALIGN.CENTER) if False else None
        tb = slide.shapes.add_textbox(Cm(x + 0.1), Cm(y0 + 0.42), Cm(4.05), Cm(0.75))
        set_text(tb, txt, size=10.5, bold=True, color=col, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(33, Cm(x + 4.35), Cm(y0 + 0.55), Cm(0.55), Cm(0.45))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MUTED
            arrow.line.color.rgb = MUTED
    add_card(slide, 2.0, 7.6, 9.2, 4.3, "DA3-only", "用 DA3 结构信号替代训练阶段 LiDAR depth supervision，检验无 LiDAR 监督是否可行。", GREEN)
    add_card(slide, 12.2, 7.6, 9.2, 4.3, "DA3+Feedback", "周期性识别风险区域和 responsible groups，生成 softpatch 信号，但不直接修改 Gaussian 参数。", ORANGE)
    add_card(slide, 22.4, 7.6, 9.2, 4.3, "PV-C", "进一步去除 LiDAR 初始化，用 COLMAP 背景 + random-box object init 维持车辆建模。", BLUE)
    notes.append("4. 这一页用流程图解释方法，不展开公式。")

    # 5
    slide = base_slide(prs, 5, "主结果：held-out test full-image 指标", "final checkpoint, test split = 245 views")
    add_table(slide, 1.3, 3.0, 30.9, 6.5, metric_rows, [
        ("label", "方法"), ("psnr", "PSNR↑"), ("ssim", "SSIM↑"), ("lpips", "LPIPS↓"), ("l1", "L1↓"), ("lidar", "LiDAR init"), ("fb", "反馈")
    ], font_size=10)
    add_card(slide, 1.6, 10.5, 9.7, 3.6, "最强可写结果", "PV-C: PSNR 25.63, SSIM 0.8468；在当前单场景 held-out test 中与 StreetGS baseline 持平或略高。", GREEN)
    add_card(slide, 12.0, 10.5, 9.7, 3.6, "反馈结论要收敛", "C 未超过 B；feedback 当前证明的是链路可运行和可审计，不是最终 RGB 指标提升。", ORANGE)
    add_card(slide, 22.4, 10.5, 9.7, 3.6, "论文写法", "主张 reduced LiDAR dependency / pure-vision viability，而不是全面 SOTA。", BLUE)
    notes.append("5. 主表显示 PV-C 是当前最有价值的结论。")

    # 6
    slide = base_slide(prs, 6, "LiDAR 使用边界审计", "区分 LiDAR 初始化与训练阶段 LiDAR 监督")
    audit_display = []
    for r in audit_rows:
        audit_display.append({
            "label": r["label"],
            "init": "是" if r["uses_lidar_init"] == "True" else "否",
            "train": r["uses_lidar_training"] or "-",
            "selected": r["uses_lidar_selected"] or "-",
            "safe": "是" if r["claim_safe_no_lidar"] == "true" else "否",
            "source": r["init_source"],
        })
    add_table(slide, 1.3, 3.0, 30.7, 6.0, audit_display, [
        ("label", "方法"), ("init", "LiDAR 初始化"), ("train", "LiDAR 训练监督"), ("selected", "LiDAR selected pixels"), ("safe", "纯视觉声明安全"), ("source", "初始化来源")
    ], font_size=9)
    add_image(slide, RESULT_ROOT / "plots" / "initialization_audit.png", 2.2, 10.0, w=12.8)
    add_card(slide, 16.8, 10.1, 13.2, 3.7, "审计结论", "只有 PV-C 满足当前 no-LiDAR initialization + no-LiDAR supervision 条件；B/C 仍是 LiDAR-init 条件。", GREEN)
    notes.append("6. 这一页防止论文表述混乱。")

    # 7
    slide = base_slide(prs, 7, "反馈链路是否真实运行？", "C 与 PV-C 均有完整 feedback trigger 证据")
    add_image(slide, RESULT_ROOT / "plots" / "feedback_trigger_timeline.png", 1.5, 3.1, w=15.3)
    fb_display = []
    for r in feedback_rows:
        if r["label"] in ["C DA3+Feedback", "PV-C Pure-Vision"]:
            fb_display.append({
                "label": r["label"],
                "status": r["status"],
                "triggers": r["trigger_count"],
                "valid": r["valid_count"],
                "pixels": r["selected_pixels_total"],
                "groups": r["gaussian_group_total"],
            })
    add_table(slide, 17.4, 3.4, 14.3, 4.2, fb_display, [
        ("label", "方法"), ("triggers", "触发"), ("valid", "有效"), ("pixels", "选中像素"), ("groups", "Gaussian groups")
    ], font_size=8.7)
    add_card(slide, 17.4, 8.4, 14.3, 3.8, "可写结论", "feedback controller 在正式训练中稳定触发：C 为 59/59，PV-C 为 30/30。", BLUE)
    add_card(slide, 17.4, 12.8, 14.3, 2.7, "边界", "这证明链路可运行，不等价于证明 feedback 提升最终 RGB 或几何指标。", ORANGE)
    notes.append("7. 区分系统可运行与性能提升。")

    # 8
    slide = base_slide(prs, 8, "训练期诊断：PSNR / L1 曲线", "diagnostic eval；不能替代 final held-out test 主表")
    add_image(slide, RESULT_ROOT / "plots" / "full_split_psnr_mean_curve.png", 1.3, 3.0, w=14.9)
    add_image(slide, RESULT_ROOT / "plots" / "full_split_l1_mean_curve.png", 17.1, 3.0, w=14.9)
    add_card(slide, 1.7, 13.2, 14.2, 2.4, "用途", "用于观察训练动态、异常点和收敛趋势。", BLUE)
    add_card(slide, 17.6, 13.2, 14.2, 2.4, "缺口", "服务器包未包含 train_loss_trace.csv，因此 loss 曲线当前不能生成。", ORANGE)
    notes.append("8. 强调诊断曲线与最终主表的区别。")

    # 9
    slide = base_slide(prs, 9, "定性结果 I：前向视角 RGB 与 Depth 对比", "iter_030000 | cam0 | view 000053_0")
    add_image(slide, RESULT_ROOT / "figures" / "formal_rgb_depth_comparisons" / "iter_030000" / "iter_030000_cam0_000053_0_rgb_depth_comparison.jpg", 0.9, 3.0, w=31.9)
    notes.append("9. 展示四组在同一视角的 RGB 与渲染深度。")

    # 10
    slide = base_slide(prs, 10, "定性结果 II：侧向/多相机视角", "iter_030000 | cam1 | view 000039_1")
    add_image(slide, RESULT_ROOT / "figures" / "formal_rgb_depth_comparisons" / "iter_030000" / "iter_030000_cam1_000039_1_rgb_depth_comparison.jpg", 0.9, 3.0, w=31.9)
    notes.append("10. 展示多相机视角下车辆和背景保留情况。")

    # 11
    slide = base_slide(prs, 11, "定性结果 III：PV-C 的价值与风险", "iter_030000 | cam4 | view 000055_4")
    add_image(slide, RESULT_ROOT / "figures" / "formal_rgb_depth_comparisons" / "iter_030000" / "iter_030000_cam4_000055_4_rgb_depth_comparison.jpg", 0.9, 3.0, w=31.9)
    add_card(slide, 1.4, 15.1, 14.8, 2.0, "可写", "PV-C 在 RGB 上保留竞争力，且满足无 LiDAR 初始化/监督审计。", GREEN)
    add_card(slide, 17.0, 15.1, 14.8, 2.0, "谨慎", "PV-C depth 色彩分布与 LiDAR-init 组差异明显，需要 final geometry metrics 支撑几何 claim。", ORANGE)
    notes.append("11. 这一页要把 PV-C 的亮点和风险同时讲清楚。")

    # 12
    slide = base_slide(prs, 12, "当前证据能支撑什么论文结论？")
    add_card(slide, 1.4, 3.2, 9.5, 4.2, "强支撑", "纯视觉 object-aware DA3+Feedback 在当前 Waymo 场景可训练、车辆保留、RGB 指标具备竞争力。", GREEN)
    add_card(slide, 12.0, 3.2, 9.5, 4.2, "中等支撑", "DA3 替代训练阶段 LiDAR supervision 的链路可运行，但指标弱于 baseline。", BLUE)
    add_card(slide, 22.6, 3.2, 9.5, 4.2, "弱/不支撑", "DA3+Feedback 目前不能 claim 全面优于 DA3-only 或 StreetGS。", RED)
    add_bullets(slide, 2.0, 8.9, 28.5, 5.2, [
        "建议论文定位：LiDAR dependency audit + pure-vision object-aware feasibility。",
        "不要写成：feedback 显著提升最终渲染质量。",
        "不要写成：几何/深度质量已被严格证明提升。",
        "Object-region 指标应标注为 object-present subset：84/245 views。",
    ], size=14)
    notes.append("12. 这是论文 claim 边界页。")

    # 13
    slide = base_slide(prs, 13, "投稿前仍需补充的材料")
    add_card(slide, 1.5, 3.0, 14.5, 3.2, "Loss trace", "恢复或重新导出 train_loss_trace.csv，补齐 loss / DA3 loss / feedback loss 曲线。", ORANGE)
    add_card(slide, 17.0, 3.0, 14.5, 3.2, "Geometry metrics", "补 final depth / geometry quantitative metrics，支撑 depth/structure claim。", ORANGE)
    add_card(slide, 1.5, 7.1, 14.5, 3.2, "More scenes", "至少补 1-2 个 Waymo scenes，否则保持 single-scene study 表述。", BLUE)
    add_card(slide, 17.0, 7.1, 14.5, 3.2, "Method schematic", "补一张可投稿的 pipeline 图：A/B/C/PV-C、DA3、feedback、LiDAR 边界。", BLUE)
    add_card(slide, 1.5, 11.2, 30.0, 2.8, "写作建议", "当前足以形成论文初稿和组会/开题汇报；若目标高水平期刊，需要补泛化和几何证据。", GREEN)
    notes.append("13. 这页是后续工作 checklist。")

    # 14
    slide = base_slide(prs, 14, "总结：当前最稳的论文故事")
    add_bullets(slide, 2.0, 3.1, 28.6, 6.0, [
        "我们不是简单追求 SOTA，而是审计 LiDAR 在 object-aware Street Gaussian reconstruction 中的作用边界。",
        "PV-C 是当前最有论文价值的结果：无 LiDAR 初始化、无 LiDAR 监督，同时保持 held-out RGB 竞争力。",
        "C/PV-C 的 feedback controller 已稳定运行，但性能提升 claim 需要更强 ablation 和几何指标。",
        "推荐题目方向：Auditing LiDAR dependency in object-aware street Gaussian reconstruction with DA3-guided feedback。",
    ], size=16)
    add_card(slide, 2.2, 10.7, 28.2, 3.4, "一句话结论", "GeoFeedback-GS 当前结果足以支撑一篇有边界的实验/方法论文初稿；补充 loss、geometry、多场景后，可升级为更强主张。", GREEN)
    notes.append("14. 结尾回到可写论文与证据边界。")

    prs.save(OUT_PPTX)
    OUT_NOTES.write_text("\n\n".join(f"## Slide {i+1}\n{n}" for i, n in enumerate(notes)), encoding="utf-8")

    reopened = Presentation(str(OUT_PPTX))
    media_dir = OUT_PPTX.with_suffix("")
    qa = {
        "pptx": str(OUT_PPTX),
        "slides": len(reopened.slides),
        "template": str(template_path) if template_path else "",
        "speaker_notes_file": str(OUT_NOTES),
        "source_result_root": str(RESULT_ROOT),
        "checks": {
            "main_table_rows": len(main_rows),
            "audit_rows": len(audit_rows),
            "feedback_rows": len(feedback_rows),
            "comparison_images_available": len(list((RESULT_ROOT / "figures" / "formal_rgb_depth_comparisons").glob("**/*.jpg"))),
        },
        "limitations": [
            "python-pptx does not create native PowerPoint speaker notes in this fast path; notes are delivered as a companion markdown file.",
            "Loss curves are absent because train_loss_trace.csv was not present in the downloaded results.",
        ],
    }
    OUT_QA.write_text(json.dumps(qa, indent=2, ensure_ascii=False), encoding="utf-8")
    print(json.dumps(qa, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()
